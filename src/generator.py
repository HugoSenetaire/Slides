from lxml import etree
from pptx import Presentation
from pptx.util import Inches


def getXmlTree(path):
    return etree.parse(path)


def getSlides(data):
    return data.getroot().xpath('diapo')


def getSlideTitle(slide):
    return (slide.xpath('title'))[0].text.replace('\n', "").replace('\r', "").replace('\t', "")


def getSlideLayout(slide):
    return slide.get("layout")


def getSlideSubtitle(slide):
    return slide.xpath('subtitle')[0].text.replace('\n', "").replace('\r', "").replace('\t', "")


def getSlideParagraph(slide):
    return slide.xpath('paragraph')[0].text.replace('\n', "").replace('\r', "").replace('\t', "").replace('<br/>', "\n")


def getSlideBullets(slide):
    return [{'text': bullet.text.replace('\n', "").replace('\r', "").replace('\t', ""), 'level': int(bullet.get("level"))} for bullet in
            slide.xpath('bullet')]


def getPresentation(path):
    return Presentation(path)


def getPresentationLayoutDictionary(presentation, ):
    presentationLayoutDictionaryIds = {
        'title': 2,
        'title+subtitle': 0,
        'title+bullets': 1,
        'only_title': 5,
        'title+bullet+image': 8,
        'title+subtitle+image': 6,
        'image': 10,
    }
    layouts = {}
    for layout in presentationLayoutDictionaryIds:
        layouts[layout] = presentation.slide_layouts[presentationLayoutDictionaryIds[layout]]
    return layouts


def addSlide(presentation, layout):
    return presentation.slides.add_slide(layout)


def setSlideTitle(new_slide, slide_title):
    new_slide.shapes.title.text = slide_title


def addTextBox(new_slide, slide_layout):
    txBox = ""
    if slide_layout == "title+paragraph":
        left = Inches(2)
        top = Inches(3)
        width = height = Inches(8)
        txBox = new_slide.shapes.add_textbox(left, top, width, height)
    elif slide_layout == "title+paragraph+image":
        left = Inches(1)
        top = Inches(3)
        width = Inches(4)
        height = Inches(6)
        txBox = new_slide.shapes.add_textbox(left, top, width, height)
    return txBox


def addImage(new_slide, slide_layout, slide):
    pic = ""
    image_source_path = slide.xpath('image')[0].xpath('source')[0].text
    placeholder = new_slide.placeholders[13]
    pic = placeholder.insert_picture(image_source_path)
    return pic


def generatePresentation(dataPath, maskPath, export=True, exportPath="output.pptx"):
    # import the xml file
    data = getXmlTree(dataPath)

    # import the presentation mask
    presentation = getPresentation(maskPath)
    # extract layout templates
    layouts = getPresentationLayoutDictionary(presentation)

    # transform it in tree
    slides = getSlides(data)

    # iterate through slides
    for i, slide in enumerate(slides):
        # slide_XXX : slide data from data XML, ONLY REPRESENTATIVE DATA (typically string), NOT DOM !
        # new_slide_XXX : slide object added in presentation

        # get slide title
        slide_title = getSlideTitle(slide)
        # get slide layout
        slide_layout = getSlideLayout(slide)

        # title only layout
        if slide_layout == "title":
            new_slide = addSlide(presentation, layouts['title'])
            setSlideTitle(new_slide, slide_title)

        # title and subtitle layout
        elif slide_layout == "title+subtitle":
            slide_subtitle = getSlideSubtitle(slide)
            new_slide = addSlide(presentation, layouts['title+subtitle'])
            setSlideTitle(new_slide, slide_title)
            new_slide_subtitle = new_slide.placeholders[1]
            new_slide_subtitle.text = slide_subtitle

        # title and a small text paragraph layout
        elif slide_layout == "title+paragraph":
            slide_paragraph = getSlideParagraph(slide)
            new_slide = addSlide(presentation, layouts['title+bullets'])
            setSlideTitle(new_slide, slide_title)
            new_slide.shapes.placeholders[1].text_frame.text=slide_paragraph

        # title and bullet point
        elif slide_layout == "title+bullets":
            slide_bullets = getSlideBullets(slide)
            new_slide = addSlide(presentation, layouts['title+bullets'])
            setSlideTitle(new_slide, slide_title)
            new_slide_bullets = new_slide.shapes.placeholders[1]
            for i, bullet in enumerate(slide_bullets):
                if i == 0:
                    new_slide_bullet_paragraph = new_slide_bullets.text_frame
                else:
                    new_slide_bullet_paragraph = new_slide_bullets.text_frame.add_paragraph()
                new_slide_bullet_paragraph.text = bullet['text']
                new_slide_bullet_paragraph.level = bullet['level']

        # title and bullet and one image
        elif slide_layout == "title+paragraph+image":
            slide_paragraph = getSlideParagraph(slide)
            new_slide = addSlide(presentation, layouts['title+bullet+image'])
            setSlideTitle(new_slide, slide_title)
            new_slide.shapes.placeholders[1].text_frame.text=slide_paragraph
            addImage(new_slide, "title+paragraph+image", slide)
        
        # juste one image 
        elif slide_layout == "image":
            new_slide = addSlide(presentation, layouts['image'])
            addImage(new_slide, "image", slide)
        
        # title and subtitle and one image
        elif slide_layout == "title+subtitle+image":
            slide_subtitle = getSlideSubtitle(slide)
            new_slide = addSlide(presentation, layouts['title+subtitle+image'])
            setSlideTitle(new_slide, slide_title)
            new_slide_subtitle = new_slide.placeholders[1]
            new_slide_subtitle.text = slide_subtitle
            addImage(new_slide, "title+image", slide)

    if export:
        presentation.save(exportPath)

    return presentation
